import os
import json
import io
import fitz  # PyMuPDF
import nbformat
import pylatexenc.latex2text
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from typing import Union, List, Optional
from markdown import markdown


class FileReaderToolkit:
    """
    A toolkit class to read and extract content from various file types, including PDFs, DOCX,
    XLSX, PPTX, HTML, Python, JSON, text files, Jupyter notebooks, LaTeX files, and Markdown files.

    This class does not instantiate any file. Each method requires a file parameter, which can be
    either a file stream or a file path.

    Attributes:
    ----------
    current_content : Optional[str]
        The current extracted content of the file.
    current_chunks : Optional[List[str]]
        The current chunked content of the file.
    current_file_type : Optional[str]
        The file type of the current file being processed.
    """

    def __init__(self):
        """Initializes the FileReaderToolkit with no file."""
        self._current_content: Optional[str] = None
        self._current_chunks: Optional[List[str]] = None
        self._current_file_type: Optional[str] = None

    @property
    def current_content(self) -> Optional[str]:
        """Gets the current extracted content of the file."""
        return self._current_content

    @property
    def current_chunks(self) -> Optional[List[str]]:
        """Gets the current chunked content of the file."""
        return self._current_chunks

    @property
    def current_file_type(self) -> Optional[str]:
        """Gets the file type of the current file being processed."""
        return self._current_file_type

    def _detect_file_type(self, file: Union[str, io.IOBase]) -> str:
        """
        Detects the file type by extracting the file extension from either the file path or
        the file-like object.

        Args:
            file (str or file-like object): The file path or file-like object.

        Returns:
            str: The file extension (e.g., '.pdf', '.docx').

        Raises:
            ValueError: If the file type cannot be determined.
        """
        if hasattr(file, "name") and isinstance(file.name, str):
            _, ext = os.path.splitext(file.name)
            if ext:
                self._current_file_type = ext.lower()
                return self._current_file_type
        elif isinstance(file, str):
            _, ext = os.path.splitext(file)
            if ext:
                self._current_file_type = ext.lower()
                return self._current_file_type
        raise ValueError(
            "Unable to detect file type: file must have a valid name or path with an extension."
        )

    def _read_file(self, file: Union[str, io.IOBase]) -> str:
        """
        Reads the content from the file based on its file type.

        Args:
            file (str or file-like object): The file path or file-like object.

        Returns:
            str: The content of the file.

        Raises:
            ValueError: If the file type is unsupported.
        """
        file_type = self._detect_file_type(file)
        file_type_handlers = {
            ".pdf": self._read_pdf,
            ".docx": self._read_docx,
            ".xlsx": self._read_xlsx,
            ".ppt": self._read_ppt,
            ".pptx": self._read_ppt,
            ".html": self._read_html,
            ".py": self._read_python,
            ".json": self._read_json,
            ".txt": self._read_text,
            ".ipynb": self._read_nb,
            ".tex": self._read_latex,
            ".md": self._read_markdown,
        }

        handler = file_type_handlers.get(file_type)
        if handler:
            self._current_content = handler(file)
            return self._current_content
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    def _read_pdf(self, file: Union[str, io.IOBase]) -> str:
        content = []
        try:
            if hasattr(file, "read"):
                file.seek(0)
                with fitz.open(stream=file.read(), filetype="pdf") as doc:
                    for page in doc:
                        content.append(page.get_text())
            else:
                with fitz.open(file) as doc:
                    for page in doc:
                        content.append(page.get_text())
        except Exception as e:
            raise ValueError(f"Error reading PDF file: {e}")
        return "\n".join(content)

    def _read_docx(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                doc = Document(file)
            else:
                doc = Document(file)
            return "\n".join(para.text for para in doc.paragraphs)
        except Exception as e:
            raise ValueError(f"Error reading DOCX file: {e}")

    def _read_xlsx(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                workbook = load_workbook(filename=file, read_only=True, data_only=True)
            else:
                workbook = load_workbook(filename=file, read_only=True, data_only=True)
            content = []
            for sheet in workbook:
                for row in sheet.iter_rows(values_only=True):
                    content.append(
                        " ".join(str(cell) for cell in row if cell is not None)
                    )
            return "\n".join(content)
        except Exception as e:
            raise ValueError(f"Error reading XLSX file: {e}")

    def _read_ppt(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                prs = Presentation(file)
            else:
                prs = Presentation(file)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content.append(shape.text)
            return "\n".join(content)
        except Exception as e:
            raise ValueError(f"Error reading PPT/PPTX file: {e}")

    def _read_html(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode("utf-8", errors="ignore")
                soup = BeautifulSoup(content, "html.parser")
            else:
                with open(file, "rb") as f:
                    content = f.read().decode("utf-8", errors="ignore")
                    soup = BeautifulSoup(content, "html.parser")
            return soup.get_text()
        except Exception as e:
            raise ValueError(f"Error reading HTML file: {e}")

    def _read_python(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode("utf-8", errors="ignore")
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            return content
        except Exception as e:
            raise ValueError(f"Error reading Python file: {e}")

    def _read_json(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                if isinstance(file, io.TextIOBase):
                    data = json.load(file)
                else:
                    content = file.read().decode("utf-8", errors="ignore")
                    data = json.loads(content)
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    data = json.load(f)
            return json.dumps(data, indent=4)
        except Exception as e:
            raise ValueError(f"Error reading JSON file: {e}")

    def _read_text(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode("utf-8", errors="ignore")
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            return content
        except Exception as e:
            raise ValueError(f"Error reading text file: {e}")

    def _read_nb(self, file: Union[str, io.IOBase]) -> str:
        try:
            if hasattr(file, "read"):
                file.seek(0)
                if isinstance(file, io.TextIOBase):
                    notebook_content = nbformat.read(file, as_version=4)
                else:
                    content = file.read().decode("utf-8", errors="ignore")
                    notebook_content = nbformat.reads(content, as_version=4)
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    notebook_content = nbformat.read(f, as_version=4)
            return nbformat.writes(notebook_content)
        except Exception as e:
            raise ValueError(f"Error reading Jupyter Notebook file: {e}")

    def _read_latex(self, file: Union[str, io.IOBase]) -> str:
        """
        Reads the content from a LaTeX file, converting LaTeX markup to plain text.

        Args:
            file (str or file-like object): The LaTeX file path or file-like object.

        Returns:
            str: The plain text content extracted from the LaTeX file.

        Raises:
            ValueError: If there's an error reading the LaTeX file.
        """
        try:
            if hasattr(file, "read"):
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode("utf-8", errors="ignore")
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            # Convert LaTeX to plain text
            latex_converter = pylatexenc.latex2text.LatexNodes2Text()
            plain_text = latex_converter.latex_to_text(content)
            return plain_text
        except Exception as e:
            raise ValueError(f"Error reading LaTeX file: {e}")

    def _read_markdown(self, file: Union[str, io.IOBase]) -> str:
        """
        Reads the content from a Markdown file, converting it to plain text.

        Args:
            file (str or file-like object): The Markdown file path or file-like object.

        Returns:
            str: The plain text content extracted from the Markdown file.

        Raises:
            ValueError: If there's an error reading the Markdown file.
        """
        try:
            if hasattr(file, "read"):
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    content = content.decode("utf-8", errors="ignore")
            else:
                with open(file, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
            # Convert Markdown to HTML, then extract text
            html_content = markdown(content)
            soup = BeautifulSoup(html_content, "html.parser")
            plain_text = soup.get_text()
            return plain_text
        except Exception as e:
            raise ValueError(f"Error reading Markdown file: {e}")

    def get_content(self, file: Union[str, io.IOBase]) -> str:
        """
        Returns the content of the file.

        Args:
            file (str or file-like object): The file path or file-like object.

        Returns:
            str: The content of the file.
        """
        return self._read_file(file)

    def chunk_content(self, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        Breaks the current content into chunks of a specified size with overlap.

        Args:
            chunk_size (int): The size of each chunk.
            chunk_overlap (int): The amount of overlap between consecutive chunks.

        Returns:
            List[str]: A list of content chunks.

        Raises:
            ValueError: If the current content is None or if parameters are invalid.
        """
        if self._current_content is None:
            raise ValueError("No content available to chunk. Please read a file first.")
        if chunk_size <= 0:
            raise ValueError("chunk_size must be a positive integer.")
        if chunk_overlap < 0:
            raise ValueError("chunk_overlap must be a non-negative integer.")
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be less than chunk_size.")

        chunks = []
        start = 0
        content_length = len(self._current_content)
        while start < content_length:
            end = min(start + chunk_size, content_length)
            chunks.append(self._current_content[start:end])
            start = end - chunk_overlap
        self._current_chunks = chunks
        return chunks
