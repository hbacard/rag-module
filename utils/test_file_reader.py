import unittest
import os
from fpdf import FPDF
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
import json
import nbformat as nbf
from file_reader_toolkit import FileReaderToolkit


# Utility functions to create sample files
def create_sample_pdf(file_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="This is a sample PDF document.", ln=True)
    pdf.output(file_path)


def create_sample_docx(file_path):
    doc = Document()
    doc.add_paragraph("This is a sample DOCX document.")
    doc.save(file_path)


def create_sample_xlsx(file_path):
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Sample"
    ws["A2"] = "XLSX"
    ws["A3"] = "File"
    wb.save(file_path)


def create_sample_pptx(file_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Sample PPTX"
    prs.save(file_path)


def create_sample_html(file_path):
    html_content = "<html><body><h1>This is a sample HTML file</h1></body></html>"
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(html_content)


def create_sample_python(file_path):
    python_code = """
def sample_function():
    return "This is a sample Python file"
"""
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(python_code)


def create_sample_json(file_path):
    data = {"name": "Sample", "type": "JSON"}
    with open(file_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4)


def create_sample_text(file_path):
    text_content = "This is a sample text file."
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text_content)


def create_sample_notebook(file_path):
    nb = nbf.v4.new_notebook()
    nb["cells"] = [nbf.v4.new_code_cell("print('This is a sample Jupyter notebook')")]
    with open(file_path, "w", encoding="utf-8") as f:
        nbf.write(nb, f)


def create_sample_latex(file_path):
    latex_content = r"""
\documentclass{article}
\begin{document}
This is a sample LaTeX document.
\end{document}
"""
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(latex_content)


def create_sample_markdown(file_path):
    markdown_content = "# Sample Markdown\n\nThis is a sample Markdown file."
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(markdown_content)


# Unit tests for the FileReaderToolkit class
class TestFileReaderToolkit(unittest.TestCase):

    @classmethod
    def setUpClass(cls):
        """Create sample files for testing."""
        os.makedirs("test_files", exist_ok=True)
        create_sample_pdf("test_files/sample.pdf")
        create_sample_docx("test_files/sample.docx")
        create_sample_xlsx("test_files/sample.xlsx")
        create_sample_pptx("test_files/sample.pptx")
        create_sample_html("test_files/sample.html")
        create_sample_python("test_files/sample.py")
        create_sample_json("test_files/sample.json")
        create_sample_text("test_files/sample.txt")
        create_sample_notebook("test_files/sample.ipynb")
        create_sample_latex("test_files/sample.tex")
        create_sample_markdown("test_files/sample.md")

    def setUp(self):
        """Initialize FileReaderToolkit before each test."""
        self.reader = FileReaderToolkit()

    def test_read_pdf(self):
        content = self.reader.get_content("test_files/sample.pdf")
        self.assertIn("This is a sample PDF document.", content)

    def test_read_docx(self):
        content = self.reader.get_content("test_files/sample.docx")
        self.assertIn("This is a sample DOCX document.", content)

    def test_read_xlsx(self):
        content = self.reader.get_content("test_files/sample.xlsx")
        self.assertIn("Sample", content)
        self.assertIn("XLSX", content)
        self.assertIn("File", content)

    def test_read_pptx(self):
        content = self.reader.get_content("test_files/sample.pptx")
        self.assertIn("Sample PPTX", content)

    def test_read_html(self):
        content = self.reader.get_content("test_files/sample.html")
        self.assertIn("This is a sample HTML file", content)

    def test_read_python(self):
        content = self.reader.get_content("test_files/sample.py")
        self.assertIn("sample_function", content)

    def test_read_json(self):
        content = self.reader.get_content("test_files/sample.json")
        self.assertIn('"name": "Sample"', content)

    def test_read_text(self):
        content = self.reader.get_content("test_files/sample.txt")
        self.assertIn("This is a sample text file.", content)

    def test_read_notebook(self):
        content = self.reader.get_content("test_files/sample.ipynb")
        self.assertIn("This is a sample Jupyter notebook", content)

    def test_read_latex(self):
        content = self.reader.get_content("test_files/sample.tex")
        self.assertIn("This is a sample LaTeX document.", content)

    def test_read_markdown(self):
        content = self.reader.get_content("test_files/sample.md")
        self.assertIn("This is a sample Markdown file.", content)

    @classmethod
    def tearDownClass(cls):
        """Remove test files after tests."""
        for file_name in os.listdir("test_files"):
            os.remove(os.path.join("test_files", file_name))
        os.rmdir("test_files")


if __name__ == "__main__":
    unittest.main()
